"""Extract slide text and speaker notes from every .pptx in a source folder.

Writes one markdown file per deck into <out_dir>/<module>.md with sections per slide,
preserving slide order, body text, and speaker notes. Used as the raw source we
synthesise the new Slides.html / Notes.md / Frameworks / Glossary / Pre-Read from.

Usage:
    python3 extract_pptx.py <source_folder> [<output_folder>]

If <output_folder> is omitted, output goes to ./_out next to this script.
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Install requirements first: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


def slide_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            line = "".join(run.text for run in para.runs).strip()
            if line:
                parts.append(line)
    return "\n".join(parts)


def notes_text(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    notes = slide.notes_slide.notes_text_frame
    if not notes:
        return ""
    return notes.text.strip()


def slugify(name: str) -> str:
    s = re.sub(r"[^A-Za-z0-9]+", "_", name).strip("_").lower()
    return s[:60]


def dump_deck(pptx_path: Path, out_dir: Path) -> Path:
    prs = Presentation(pptx_path)
    out = out_dir / f"{slugify(pptx_path.stem)}.md"
    with out.open("w") as f:
        f.write(f"# {pptx_path.name}\n\n")
        for idx, slide in enumerate(prs.slides, start=1):
            body = slide_text(slide)
            notes = notes_text(slide)
            f.write(f"\n---\n\n## Slide {idx}\n\n")
            if body:
                f.write("### Body\n\n")
                f.write(body + "\n\n")
            if notes:
                f.write("### Speaker Notes\n\n")
                f.write(notes + "\n\n")
    return out


def main() -> None:
    if len(sys.argv) < 2:
        print(__doc__, file=sys.stderr)
        sys.exit(1)
    source = Path(sys.argv[1]).expanduser().resolve()
    out_dir = Path(sys.argv[2]).expanduser().resolve() if len(sys.argv) >= 3 else Path(__file__).resolve().parent / "_out"
    out_dir.mkdir(parents=True, exist_ok=True)
    if not source.exists():
        print(f"Source folder not found: {source}", file=sys.stderr)
        sys.exit(1)
    decks = sorted(source.rglob("*.pptx"))
    if not decks:
        print("No .pptx files found.", file=sys.stderr)
        sys.exit(1)
    for deck in decks:
        out = dump_deck(deck, out_dir)
        print(f"wrote {out}  ({deck})")


if __name__ == "__main__":
    main()
